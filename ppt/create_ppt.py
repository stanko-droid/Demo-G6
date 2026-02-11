#!/usr/bin/env python3
"""
create_ppt.py
Genererar en PowerPoint från Markdown + manuella speaker notes
Kör: python create_ppt.py
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm

MD_FILE = "SLUT_LOGG.md"
OUTPUT_FILE = "NewsFlash_MasterClass.pptx"


# -------------------------------------------------
# 1. PARSA MARKDOWN
# -------------------------------------------------
def parse_markdown(md_text):
    slides = []
    current_slide = None
    in_code_block = False
    code_buffer = []
    diagram_buffer = []

    for line in md_text.split("\n"):
        # Ny slide vid H1/H2/H3
        if re.match(r"^#{1,3} ", line):
            if current_slide:
                slides.append(current_slide)

            title = re.sub(r"^#{1,3} ", "", line).strip()
            current_slide = {
                "title": title,
                "bullets": [],
                "code": [],
                "diagram": [],
            }
            continue

        # Kodblock start/slut
        if line.strip().startswith("```"):
            if not in_code_block:
                in_code_block = True
                code_buffer = []
            else:
                in_code_block = False
                if current_slide:
                    current_slide["code"].append("\n".join(code_buffer))
            continue

        if in_code_block:
            code_buffer.append(line)
            continue

        # Mermaid diagram detection
        if "graph " in line or "sequenceDiagram" in line:
            if current_slide:
                current_slide["diagram"].append("Mermaid diagram detected – insert exported image here.")
            continue

        # Bullet points
        if line.strip().startswith("- ") or line.strip().startswith("•"):
            if current_slide:
                bullet = line.strip()[2:]
                current_slide["bullets"].append(bullet)
            continue

        # Vanlig text som kort punkt
        if current_slide and line.strip():
            current_slide["bullets"].append(line.strip())

    if current_slide:
        slides.append(current_slide)

    return slides


# -------------------------------------------------
# 2. SPEAKER NOTES PER SLIDE
# -------------------------------------------------
def generate_speaker_notes(title):
    if "Engineering Culture" in title:
        return "Talare: Sofia\n\nFokusera på kultur, Jira och 'No Ticket No Code'.\nTidsram: 10–12 minuter.\nPoängtera spårbarhet och struktur."
    if "Architecture" in title or "3-Tier" in title:
        return "Talare: Ludwig\n\nFörklara separation of concerns.\nPresentation / Business / Data.\nTidsram: 10–12 minuter."
    if "Security" in title:
        return "Talare: Aisa\n\nDefense in Depth, OWASP, OIDC.\nFörklara varför OIDC är säkrare än secrets.\nTidsram: 10–12 minuter."
    if "Test" in title:
        return "Talare: Stanko\n\n90 tester. Integration vs unit.\nNormalization bug story.\nTidsram: 10–12 minuter."
    if "CI" in title or "DevOps" in title:
        return "Talare: Hashim\n\n13 runs. OIDC-debugging.\nPipeline steg för steg.\nTidsram: 10–12 minuter."
    return "Talare: Team\n\nFölj manus och håll slide-text kort."


# -------------------------------------------------
# 3. SKAPA PPT
# -------------------------------------------------
def create_presentation(slides):
    prs = Presentation()

    for slide_data in slides:
        # Title + Content layout
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # Titel
        slide.shapes.title.text = slide_data["title"]

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        tf.word_wrap = True

        # Bullets
        for i, bullet in enumerate(slide_data["bullets"][:6]):  # max 6 per slide
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 1

        # Kod-platshållare
        if slide_data["code"]:
            left = Inches(0.5)
            top = Inches(4.5)
            width = Inches(9)
            height = Inches(1.5)
            box = slide.shapes.add_textbox(left, top, width, height)
            box.text_frame.text = "[CODE PLACEHOLDER]\n" + slide_data["code"][0]
            for p in box.text_frame.paragraphs:
                p.font.name = "Courier New"
                p.font.size = Pt(12)

        # Diagram-platshållare
        if slide_data["diagram"]:
            left = Inches(0.5)
            top = Inches(4)
            width = Inches(9)
            height = Inches(1)
            box = slide.shapes.add_textbox(left, top, width, height)
            box.text_frame.text = "[DIAGRAM PLACEHOLDER]\nExport Mermaid → PNG och infoga här."

        # Speaker Notes
        notes = slide.notes_slide.notes_text_frame
        notes.text = generate_speaker_notes(slide_data["title"])

    prs.save(OUTPUT_FILE)


# -------------------------------------------------
# MAIN
# -------------------------------------------------
def main():
    with open(MD_FILE, "r", encoding="utf-8") as f:
        md_text = f.read()

    slides = parse_markdown(md_text)
    create_presentation(slides)
    print(f"Presentation skapad: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()